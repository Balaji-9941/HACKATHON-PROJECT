import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Theme Colors
    BG_DARK = RGBColor(8, 12, 22)        # #080c16
    CARD_DARK = RGBColor(15, 23, 42)      # #0f172a
    CARD_BORDER = RGBColor(30, 41, 59)    # #1e293b
    NEON_BLUE = RGBColor(59, 130, 246)    # #3b82f6
    ROYAL_BLUE = RGBColor(37, 99, 235)    # #2563eb
    EMERALD = RGBColor(16, 185, 129)      # #10b981
    CRIMSON = RGBColor(239, 68, 68)       # #ef4444
    AMBER = RGBColor(245, 158, 11)        # #f59e0b
    WHITE = RGBColor(248, 250, 252)       # #f8fafc
    SLATE_LIGHT = RGBColor(203, 213, 225) # #cbd5e1
    SLATE_MUTED = RGBColor(148, 163, 184) # #94a3b8

    slide_layout = prs.slide_layouts[6] # Blank layout

    def apply_dark_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_DARK
        bg.line.color.rgb = BG_DARK
        return bg

    def add_header(slide, title_text, subtitle_text="", slide_num=""):
        # Top accent bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.04))
        bar.fill.solid()
        bar.fill.fore_color.rgb = ROYAL_BLUE
        bar.line.color.rgb = ROYAL_BLUE

        # Title box
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(9.5), Inches(0.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(21)
        p.font.bold = True
        p.font.color.rgb = WHITE

        if subtitle_text:
            p2 = tf.add_paragraph()
            p2.text = subtitle_text
            p2.font.size = Pt(11)
            p2.font.color.rgb = SLATE_MUTED

        # Watermark Right
        txR = slide.shapes.add_textbox(Inches(9.8), Inches(0.5), Inches(2.7), Inches(0.8))
        tf_r = txR.text_frame
        p_r = tf_r.paragraphs[0]
        p_r.text = "Team Chronoforge | Hexaware '26"
        p_r.alignment = PP_ALIGN.RIGHT
        p_r.font.size = Pt(10)
        p_r.font.bold = True
        p_r.font.color.rgb = NEON_BLUE
        if slide_num:
            p_rn = tf_r.add_paragraph()
            p_rn.text = f"Slide {slide_num} / 16"
            p_rn.alignment = PP_ALIGN.RIGHT
            p_rn.font.size = Pt(9.5)
            p_rn.font.color.rgb = SLATE_MUTED

    def add_screenshot_placeholder(slide, left, top, width, height, label="Screenshot Placeholder", aspect="16:9"):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(12, 18, 34)
        box.line.color.rgb = NEON_BLUE
        box.line.width = Pt(1.5)
        
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"📸 [{label} — {aspect}]"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = NEON_BLUE

        p2 = tf.add_paragraph()
        p2.text = "Swap with actual application export screenshot"
        p2.alignment = PP_ALIGN.CENTER
        p2.font.size = Pt(9)
        p2.font.color.rgb = SLATE_MUTED
        return box

    # ----------------------------------------------------
    # SLIDE 1: Title Slide & Executive Hook
    # ----------------------------------------------------
    s1 = prs.slides.add_slide(slide_layout)
    apply_dark_bg(s1)

    # Hero Placeholder
    add_screenshot_placeholder(s1, Inches(0.8), Inches(0.8), Inches(11.733), Inches(5.9), "Admin Command Center Hero View", "Full-Bleed Background")

    # Title Card Overlay
    t_box = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(1.2), Inches(10.9), Inches(5.1))
    t_box.fill.solid()
    t_box.fill.fore_color.rgb = RGBColor(10, 15, 29)
    t_box.line.color.rgb = ROYAL_BLUE
    t_box.line.width = Pt(2)
    tf1 = t_box.text_frame
    tf1.word_wrap = True

    p = tf1.paragraphs[0]
    p.text = "HEXAWARE MAVERICKS HACKATHON 2026 | TECHNICAL JUDGING DECK"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = NEON_BLUE

    p_t = tf1.add_paragraph()
    p_t.text = "PayTelemetry: Autonomous In-Flight Fraud Intelligence"
    p_t.font.size = Pt(30)
    p_t.font.bold = True
    p_t.font.color.rgb = WHITE

    p_sub = tf1.add_paragraph()
    p_sub.text = "Sub-20ms Continuous Anomaly Scoring, Risk-Adaptive Graduated Friction & Real-Time TreeSHAP Explainability"
    p_sub.font.size = Pt(14)
    p_sub.font.color.rgb = SLATE_MUTED

    tf1.add_paragraph()
    p_b = tf1.add_paragraph()
    p_b.text = "★ 100.0% Precision | 99.87% Recall | Sub-20ms Tier 1 Latency | Zero Customer Block Churn ★"
    p_b.font.size = Pt(13)
    p_b.font.bold = True
    p_b.font.color.rgb = EMERALD

    tf1.add_paragraph()
    p_team = tf1.add_paragraph()
    p_team.text = "TEAM CHRONOFORGE:"
    p_team.font.size = Pt(12)
    p_team.font.bold = True
    p_team.font.color.rgb = WHITE

    members = [
        "1. BALAJI M — Team Lead & ML Systems Architect (XGBoost, TreeSHAP, Feature Engineering)",
        "2. DHARSHINI G S — Full-Stack & Streaming Pipeline Engineer (FastAPI Microservices, React Hot-Path)",
        "3. DAKSHANYA D — Backend & Telemetry Engineer (Node.js Gateway, Signal Extractors, Circuit Breakers)",
        "4. DEENA SHEK X — UI/UX & SOC Analytics Engineer (Investigator Console, Graph Networks, Kanban Triage)"
    ]
    for m in members:
        pm = tf1.add_paragraph()
        pm.text = f"• {m}"
        pm.font.size = Pt(10.5)
        pm.font.color.rgb = SLATE_LIGHT

    # ----------------------------------------------------
    # SLIDE 2: The Modern Payment Dilemma & Market Failure
    # ----------------------------------------------------
    s2 = prs.slides.add_slide(slide_layout)
    apply_dark_bg(s2)
    add_header(s2, "The Real-Time Payment Dilemma", "Why Legacy Rule Engines and Asynchronous Batch Models Fail Modern Instant Rails (UPI / Fast Payments)", "2")

    cards_s2 = [
        {"tag": "⚡ THE SETTLEMENT RACE", "col": CRIMSON, "title": "Irrevocable Seconds", "desc": "Real-time payment rails settle in seconds. Traditional batch ML models run post-settlement — flagging fraud only after money has irrevocably left the banking ecosystem."},
        {"tag": "📉 REVENUE DESTRUCTION", "col": AMBER, "title": "The False Decline Problem", "desc": "Rigid binary 'Block vs Allow' rule sets penalize loyal customers. Industry-wide, false declines cost merchants many times more in lost legitimate revenue than fraud itself."},
        {"tag": "🔍 BLACK-BOX PARALYSIS", "col": NEON_BLUE, "title": "SOC Alert Fatigue", "desc": "Security analysts receive ambiguous risk scores without mathematical explainability, spending significant manual time per investigation while automated mule rings drain balances."}
    ]
    for i, c in enumerate(cards_s2):
        bx = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + i*4.0), Inches(1.6), Inches(3.733), Inches(5.1))
        bx.fill.solid()
        bx.fill.fore_color.rgb = CARD_DARK
        bx.line.color.rgb = CARD_BORDER
        tf = bx.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = c["tag"]
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = c["col"]

        p_t = tf.add_paragraph()
        p_t.text = c["title"]
        p_t.font.size = Pt(16)
        p_t.font.bold = True
        p_t.font.color.rgb = WHITE

        tf.add_paragraph()
        p_d = tf.add_paragraph()
        p_d.text = c["desc"]
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = SLATE_LIGHT

    # ----------------------------------------------------
    # SLIDE 3: The Chronoforge Philosophy & Core Solution
    # ----------------------------------------------------
    s3 = prs.slides.add_slide(slide_layout)
    apply_dark_bg(s3)
    add_header(s3, "The Chronoforge Paradigm: In-Flight Adaptive Defense", "Replacing Post-Mortem Analytics with Hot-Path Prevention and Risk-Calibrated Friction", "3")

    quads_s3 = [
        {"title": "1. Pre-Debit Hot-Path Scoring", "desc": "Scoring 7 continuous telemetry signals synchronously in Node.js before settlement — no slow external network call sits on the critical blocking path.", "col": NEON_BLUE},
        {"title": "2. Dynamic Graduated Friction", "desc": "Eliminating binary hard-blocks by introducing risk-proportionate user challenges (Biometric Step-Up, In-App Confirmation).", "col": EMERALD},
        {"title": "3. Layered Explainable AI", "desc": "A deterministic rule-and-anomaly waterfall runs on every transaction; a real trained XGBoost model with TreeSHAP attribution enriches the score asynchronously when available — never blocking, never faked.", "col": AMBER},
        {"title": "4. Full-Cycle SOC Operations", "desc": "Real-time WebSocket streaming feed, interactive Kanban incident triage, multi-hop entity graph intelligence, and attack scenario injectors.", "col": CRIMSON}
    ]
    for i, q in enumerate(quads_s3):
        col = i % 2
        row = i // 2
        bx = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + col*5.966), Inches(1.6 + row*2.7), Inches(5.766), Inches(2.5))
        bx.fill.solid()
        bx.fill.fore_color.rgb = CARD_DARK
        bx.line.color.rgb = CARD_BORDER
        tf = bx.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = q["title"]
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = q["col"]

        tf.add_paragraph()
        p_d = tf.add_paragraph()
        p_d.text = q["desc"]
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = SLATE_LIGHT

    # ----------------------------------------------------
    # SLIDE 4: Architectural Blueprint & Hot-Path Pipeline
    # ----------------------------------------------------
    s4 = prs.slides.add_slide(slide_layout)
    apply_dark_bg(s4)
    add_header(s4, "Layered Microservices Architecture", "A Synchronous Tier 1 Engine That Never Fails, Enriched by an Optional Tier 2 ML Layer", "4")

    arch_layers = [
        {"title": "📱 CLIENT LAYER (React / Vite)", "col": NEON_BLUE, "items": ["• Real-time pre-check telemetry sensor", "• Adaptive checkout UI", "• Biometric step-up modal", "• Scenario & auto-flow trigger dashboard"]},
        {"title": "⚙️ GATEWAY & TELEMETRY (Node.js)", "col": EMERALD, "items": ["• 7-signal deterministic risk engine (Tier 1)", "• Synchronous sub-20ms execution", "• Non-blocking ML circuit breaker", "• Real-time Socket.io event broadcaster", "• MongoDB persistence & audit logs"]},
        {"title": "🤖 ML MICROSERVICE (Python FastAPI)", "col": CRIMSON, "items": ["• Balanced XGBoost Classifier v4", "• Real SHAP TreeExplainer attribution", "• Asynchronous score enrichment", "• Never blocks or delays settlement", "• Continuous drift monitoring"]}
    ]
    for i, a in enumerate(arch_layers):
        bx = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + i*4.0), Inches(1.6), Inches(3.733), Inches(5.1))
        bx.fill.solid()
        bx.fill.fore_color.rgb = CARD_DARK
        bx.line.color.rgb = CARD_BORDER
        tf = bx.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = a["title"]
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = a["col"]

        tf.add_paragraph()
        for item in a["items"]:
            pi = tf.add_paragraph()
            pi.text = item
            pi.font.size = Pt(11.5)
            pi.font.color.rgb = SLATE_LIGHT

    # ----------------------------------------------------
    # SLIDE 5: Innovation 1 — Multi-Vector Balanced XGBoost
    # ----------------------------------------------------
    s5 = prs.slides.add_slide(slide_layout)
    apply_dark_bg(s5)
    add_header(s5, "Deep Dive: Overcoming Single-Feature Bias in Fraud AI", "How We Address Account-Drain Over-Reliance in Naive Decision Trees", "5")

    b_s5_1 = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(11.733), Inches(1.5))
    b_s5_1.fill.solid()
    b_s5_1.fill.fore_color.rgb = CARD_DARK
    b_s5_1.line.color.rgb = CARD_BORDER
    tf = b_s5_1.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "⚠️ The Raw Dataset Trap:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = AMBER
    p2 = tf.add_paragraph()
    p2.text = "Financial fraud datasets are frequently dominated by full-account-drain patterns (96.5%), which causes naive tree splits to assign near-zero risk to 1500× baseline spikes or rapid velocity flurries if the customer holds a high account balance."
    p2.font.size = Pt(11.5)
    p2.font.color.rgb = SLATE_LIGHT

    b_s5_2 = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.3), Inches(5.766), Inches(3.4))
    b_s5_2.fill.solid()
    b_s5_2.fill.fore_color.rgb = CARD_DARK
    b_s5_2.line.color.rgb = CARD_BORDER
    tf = b_s5_2.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🛠️ The Chronoforge Fix:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = EMERALD
    lines_fix = [
        "• Sub-feature bagging (colsample_bytree = 0.65, colsample_bylevel = 0.65).",
        "• L1 (reg_alpha = 0.1) & L2 (reg_lambda = 1.0) regularization.",
        "• Multi-anomaly calibration across velocity flurries, high amount-ratio spikes, device takeovers, and geographic displacement."
    ]
    for lf in lines_fix:
        pl = tf.add_paragraph()
        pl.text = lf
        pl.font.size = Pt(11)
        pl.font.color.rgb = SLATE_LIGHT

    b_s5_3 = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.766), Inches(3.3), Inches(5.766), Inches(3.4))
    b_s5_3.fill.solid()
    b_s5_3.fill.fore_color.rgb = CARD_DARK
    b_s5_3.line.color.rgb = CARD_BORDER
    tf = b_s5_3.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "📊 Actual Training Result (metrics.json):"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = NEON_BLUE
    p_res = tf.add_paragraph()
    p_res.text = "Achieved 0.9998 ROC-AUC on 22,522 held-out test transactions with feature importance balanced across 5+ telemetry dimensions (Location 25.0%, Device 23.4%, Rule Synergy 20.9%, Drain 15.3%, Amount 4.7%) rather than a single dominant account-drain signal."
    p_res.font.size = Pt(11.5)
    p_res.font.color.rgb = WHITE

    # ----------------------------------------------------
    # SLIDE 6: Innovation 2 — Real-Time In-Flight TreeSHAP
    # ----------------------------------------------------
    s6 = prs.slides.add_slide(slide_layout)
    apply_dark_bg(s6)
    add_header(s6, "Deep Dive: Explainable AI in the Scoring Path", "Converting Model Probabilities into Real Mathematical Attribution — Never Hand-Assigned", "6")

    b_s6_l = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(6.8), Inches(5.1))
    b_s6_l.fill.solid()
    b_s6_l.fill.fore_color.rgb = CARD_DARK
    b_s6_l.line.color.rgb = CARD_BORDER
    tf = b_s6_l.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "How It Works:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = NEON_BLUE

    s6_bullets = [
        "• Exact Shapley Decomposition: Computes exact additive feature contributions via shap.TreeExplainer on the actual trained model: Score = base + sum(phi_i).",
        "• Real Waterfall Breakdown: Baseline Deviation (+50 pts), Velocity Burst (+30 pts), Unrecognized Hardware (+15 pts).",
        "• Dual-Audience Natural Language Engine:",
        "    - Consumers: 'Payment of ₹10,00,000 is 1538.5× your average transaction (₹650).'",
        "    - Investigators: Full multi-factor telemetry audit log with timestamped device fingerprints.",
        "• Sub-Millisecond Speed: TreeSHAP executes in low single-digit ms on the compiled model asynchronously."
    ]
    for b in s6_bullets:
        pb = tf.add_paragraph()
        pb.text = b
        pb.font.size = Pt(11)
        pb.font.color.rgb = SLATE_LIGHT

    add_screenshot_placeholder(s6, Inches(7.8), Inches(1.6), Inches(4.733), Inches(5.1), "TelemetryDrawer / ScoreBreakdownChart Waterfall", "4:3")

    # ----------------------------------------------------
    # SLIDE 7: Innovation 3 — 5-Stage Dynamic Graduated Friction
    # ----------------------------------------------------
    s7 = prs.slides.add_slide(slide_layout)
    apply_dark_bg(s7)
    add_header(s7, "Deep Dive: Risk-Adaptive UX vs Binary Blocking", "Delivering Friction-Free Checkout for Legitimate Users and Targeted Friction for Anomalies", "7")

    b_s7_t = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.733), Inches(1.7))
    b_s7_t.fill.solid()
    b_s7_t.fill.fore_color.rgb = CARD_DARK
    b_s7_t.line.color.rgb = CARD_BORDER
    tf = b_s7_t.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "The 5 Graduated Friction Tiers:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = EMERALD
    f_tiers = [
        "• Tier 0 (0–30 | NONE): Zero friction (<15ms)  |  • Tier 1 (31–50 | BANNER): Security warning banner",
        "• Tier 2 (51–70 | CONFIRM): Explicit review modal  |  • Tier 3 (71–85 | STEP-UP): Biometric Challenge",
        "• Tier 4 (86–100 | STEP-UP + ALERT): High-friction challenge + Automated SOC Incident"
    ]
    for ft in f_tiers:
        p_ft = tf.add_paragraph()
        p_ft.text = ft
        p_ft.font.size = Pt(11)
        p_ft.font.color.rgb = WHITE

    # 3 Placeholder frames side-by-side
    add_screenshot_placeholder(s7, Inches(0.8), Inches(3.4), Inches(3.644), Inches(3.3), "PaymentReview Normal Flow", "9:19.5")
    add_screenshot_placeholder(s7, Inches(4.844), Inches(3.4), Inches(3.644), Inches(3.3), "RiskWarningModal Tier 2", "9:19.5")
    add_screenshot_placeholder(s7, Inches(8.888), Inches(3.4), Inches(3.644), Inches(3.3), "StepUpAuth Biometric Tier 3/4", "9:19.5")

    # ----------------------------------------------------
    # SLIDE 8: Innovation 4 — Entity Graph & Mule Network Intelligence
    # ----------------------------------------------------
    s8 = prs.slides.add_slide(slide_layout)
    apply_dark_bg(s8)
    add_header(s8, "Graph Telemetry: Catching Coordinated Mule Rings", "Identifying Multi-Hop Syndicates and High-Velocity Layering Networks", "8")

    b_s8_t = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.733), Inches(1.8))
    b_s8_t.fill.solid()
    b_s8_t.fill.fore_color.rgb = CARD_DARK
    b_s8_t.line.color.rgb = CARD_BORDER
    tf = b_s8_t.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Core Graph Capabilities:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = NEON_BLUE
    g_caps = [
        "• Network Topology Scoring: Tracks velocity of money movement across shared hardware keys and clusters.",
        "• Multi-Hop Fan-In / Fan-Out: Flags aggregator accounts consolidating dispersed small payments into crypto/wire exits.",
        "• Proximity Penalty: Automatically inflates risk scores for unflagged accounts within 2 hops of flagged nodes."
    ]
    for gc in g_caps:
        p_gc = tf.add_paragraph()
        p_gc.text = gc
        p_gc.font.size = Pt(11)
        p_gc.font.color.rgb = SLATE_LIGHT

    add_screenshot_placeholder(s8, Inches(0.8), Inches(3.5), Inches(11.733), Inches(3.2), "FraudNetworkGraph.jsx Live Force-Directed Mule Cluster", "16:9")

    # ----------------------------------------------------
    # SLIDE 9: Fault Tolerance, Resilience & Circuit Breaking
    # ----------------------------------------------------
    s9 = prs.slides.add_slide(slide_layout)
    apply_dark_bg(s9)
    add_header(s9, "Enterprise Reliability: The ML Layer Can Never Take Down Payments", "How PayTelemetry Guarantees Scoring Continuity Even During Service Disruptions", "9")

    b_s9_l = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(6.8), Inches(5.1))
    b_s9_l.fill.solid()
    b_s9_l.fill.fore_color.rgb = CARD_DARK
    b_s9_l.line.color.rgb = CARD_BORDER
    tf = b_s9_l.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Resilience Architecture:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = EMERALD

    res_items = [
        "1. Non-Blocking Circuit Breaker: Monitors Python ML microservice with a 150ms timeout. After 3 failures, circuit trips OPEN immediately.",
        "2. Tier 1 Is the System of Record: Node.js gateway scores every transaction synchronously & independently (<20ms). If Tier 2 is down, Tier 1 produces valid real-time scores.",
        "3. Auto-Healing & Half-Open Probing: Automatically probes Python service every 10s; restores Tier 2 on recovery with zero manual intervention.",
        "4. Immutable Audit Ledger: Every transaction, score, model tier used, and investigator action is logged to MongoDB."
    ]
    for r in res_items:
        pr = tf.add_paragraph()
        pr.text = r
        pr.font.size = Pt(11)
        pr.font.color.rgb = SLATE_LIGHT

    add_screenshot_placeholder(s9, Inches(7.8), Inches(1.6), Inches(4.733), Inches(5.1), "SystemHealthPanel.jsx Live Service Status Dots", "4:3")

    # ----------------------------------------------------
    # SLIDE 10: Empirical Results & Performance Benchmarks
    # ----------------------------------------------------
    s10 = prs.slides.add_slide(slide_layout)
    apply_dark_bg(s10)
    add_header(s10, "Validation on Real, Held-Out Transaction Data", "Measured Results — Not Projected (Tested on 22,522 Held-Out Rows from 112,607 Dataset)", "10")

    # 4 Stat Cards Top
    metrics_s10 = [
        {"val": "100.0%", "label": "PRECISION (0 False Positives)", "col": EMERALD},
        {"val": "99.87%", "label": "RECALL (Sensitivity)", "col": NEON_BLUE},
        {"val": "0.9998", "label": "ROC-AUC DISCRIMINATION", "col": CRIMSON},
        {"val": "<20ms", "label": "TIER 1 SCORING LATENCY", "col": WHITE}
    ]
    for i, m in enumerate(metrics_s10):
        bx = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + i*3.0), Inches(1.5), Inches(2.733), Inches(1.5))
        bx.fill.solid()
        bx.fill.fore_color.rgb = CARD_DARK
        bx.line.color.rgb = CARD_BORDER
        tf = bx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = m["val"]
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = m["col"]
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = m["label"]
        p2.font.size = Pt(9.5)
        p2.font.bold = True
        p2.font.color.rgb = SLATE_MUTED
        p2.alignment = PP_ALIGN.CENTER

    # Left: Confusion Matrix & Features
    b_s10_l = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.2), Inches(6.0), Inches(3.5))
    b_s10_l.fill.solid()
    b_s10_l.fill.fore_color.rgb = CARD_DARK
    b_s10_l.line.color.rgb = CARD_BORDER
    tf = b_s10_l.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Held-Out Confusion Matrix & Feature Weights:"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = WHITE

    cm_txts = [
        "• True Negatives (TN): 19,462  |  True Positives (TP): 3,056",
        "• False Positives (FP): 0 (Zero customer harassment!)",
        "• False Negatives (FN): 4 (99.87% capture rate)",
        "• Top Features: Location (25.0%), Device (23.4%), Rule Synergy (20.9%), Account Drain (15.3%), Amount Ratio (4.7%)."
    ]
    for cmt in cm_txts:
        pcm = tf.add_paragraph()
        pcm.text = cmt
        pcm.font.size = Pt(10.5)
        pcm.font.color.rgb = SLATE_LIGHT

    add_screenshot_placeholder(s10, Inches(7.0), Inches(3.2), Inches(5.533), Inches(3.5), "ModelPerformance.jsx Real Confusion Matrix & ROC Curve", "4:3")

    # ----------------------------------------------------
    # SLIDE 11: Validated Attack Scenarios (Live System Demo)
    # ----------------------------------------------------
    s11 = prs.slides.add_slide(slide_layout)
    apply_dark_bg(s11)
    add_header(s11, "Live Attack Vector Injections & Scoring Behavior", "Verified Behavioral Accuracy Across 5 Critical Financial Attack Scenarios", "11")

    b_s11_l = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(6.5), Inches(5.2))
    b_s11_l.fill.solid()
    b_s11_l.fill.fore_color.rgb = CARD_DARK
    b_s11_l.line.color.rgb = CARD_BORDER
    tf = b_s11_l.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "5 Verified Attack Scenarios:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = NEON_BLUE

    sc_list = [
        "1. Account Takeover / Drain: Large balance extraction from novel device → CRITICAL → Biometric Step-Up + SOC Incident.",
        "2. Coordinated Mule Ring: Rapid fund layering into crypto counterparty → HIGH/CRITICAL → Graph Flag + Escalation.",
        "3. Impossible Travel Anomaly: 6,000km geo-hop to Moscow at 3:00 AM → MEDIUM → In-App Confirmation Modal.",
        "4. Extreme Baseline Multiple: ₹10,00,000 payment (1538× normal ₹650 baseline) to contact → HIGH → Biometric Step-Up.",
        "5. Legitimate Everyday Payment: ₹350 lunch payment to Swiggy from known Pixel 8 → NONE → Frictionless Pass in <15ms."
    ]
    for sc in sc_list:
        psc = tf.add_paragraph()
        psc.text = sc
        psc.font.size = Pt(10.5)
        psc.font.color.rgb = SLATE_LIGHT

    add_screenshot_placeholder(s11, Inches(7.5), Inches(1.5), Inches(5.033), Inches(5.2), "LiveStreamTable & AlertQueue Mid-Scenario", "16:9")

    # ----------------------------------------------------
    # SLIDE 12: Business Impact — Projected Value
    # ----------------------------------------------------
    s12 = prs.slides.add_slide(slide_layout)
    apply_dark_bg(s12)
    add_header(s12, "Where the Financial Value Comes From", "Directional Estimates Based on Industry Benchmarks — Not Measured Production Results", "12")

    biz_items = [
        {"title": "💰 Fraud Loss Reduction", "col": CRIMSON, "desc": "Pre-debit intervention is designed to stop account drains before settlement rather than after, which is where most recovery value is lost industry-wide."},
        {"title": "📈 GMV Preservation via Graduated Friction", "col": EMERALD, "desc": "Replacing binary blocking with tiered friction is designed to recover transactions that would otherwise be lost to over-aggressive false declines."},
        {"title": "⚡ SOC Efficiency", "col": NEON_BLUE, "desc": "Real SHAP-based explanations are designed to cut investigator review time per alert versus a black-box score, enabling higher monitored volume per analyst."}
    ]
    for i, b in enumerate(biz_items):
        bx = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + i*4.0), Inches(1.6), Inches(3.733), Inches(5.1))
        bx.fill.solid()
        bx.fill.fore_color.rgb = CARD_DARK
        bx.line.color.rgb = CARD_BORDER
        tf = bx.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = b["title"]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = b["col"]

        tf.add_paragraph()
        p_d = tf.add_paragraph()
        p_d.text = b["desc"]
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = SLATE_LIGHT

    # ----------------------------------------------------
    # SLIDE 13: Regulatory Alignment, Governance & Ethics
    # ----------------------------------------------------
    s13 = prs.slides.add_slide(slide_layout)
    apply_dark_bg(s13)
    add_header(s13, "Designed With Institutional Compliance Principles in Mind", "Architected in Alignment With RBI/NPCI, GDPR, and Responsible AI Directives — Not a Certified Compliance Claim", "13")

    reg_quads = [
        {"title": "RBI / NPCI Alignment", "desc": "Architecture supports real-time risk-based authentication patterns consistent with digital payment intermediary guidance (not a certification).", "col": NEON_BLUE},
        {"title": "Explainability by Design", "desc": "TreeSHAP attribution and the Tier 1 rule waterfall are structured to support a 'right to explanation' for automated decisions (GDPR Article 22 / DPDP Act).", "col": EMERALD},
        {"title": "Model Governance", "desc": "Tracks prediction distributions and model-tier usage over time to support future drift monitoring.", "col": AMBER},
        {"title": "Minimal PII in ML Path", "desc": "Features passed to the ML layer are numeric/normalized signals, not raw personal identifiers.", "col": WHITE}
    ]
    for i, rq in enumerate(reg_quads):
        col = i % 2
        row = i // 2
        bx = s13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + col*5.966), Inches(1.6 + row*2.7), Inches(5.766), Inches(2.5))
        bx.fill.solid()
        bx.fill.fore_color.rgb = CARD_DARK
        bx.line.color.rgb = CARD_BORDER
        tf = bx.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = rq["title"]
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = rq["col"]

        tf.add_paragraph()
        p_d = tf.add_paragraph()
        p_d.text = rq["desc"]
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = SLATE_LIGHT

    # ----------------------------------------------------
    # SLIDE 14: Competitive Matrix: Chronoforge vs The World
    # ----------------------------------------------------
    s14 = prs.slides.add_slide(slide_layout)
    apply_dark_bg(s14)
    add_header(s14, "How Chronoforge Differentiates", "Comprehensive Feature Comparison Matrix", "14")

    # Table layout
    rows = 8
    cols = 4
    left = Inches(0.8)
    top = Inches(1.6)
    width = Inches(11.733)
    height = Inches(5.2)

    table_shape = s14.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    table.columns[0].width = Inches(2.6)
    table.columns[1].width = Inches(2.8)
    table.columns[2].width = Inches(2.8)
    table.columns[3].width = Inches(3.533)

    headers = ["Capability", "Legacy Rule Engines", "Generic ML Hackathons", "PayTelemetry (Chronoforge)"]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(15, 23, 42)
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = WHITE

    matrix_rows = [
        ["Intervention Model", "Binary Block/Allow", "Binary Flagging", "5-Stage Dynamic Graduated Friction"],
        ["Scoring Hot-Path", "Fast but inflexible", "Often >200ms, asynchronous", "<20ms Tier 1, synchronous, always available"],
        ["Explainability", "Hardcoded strings", "None (Black-Box)", "Real Rule Waterfall + Real TreeSHAP (Tier 2)"],
        ["Failure Handling", "N/A", "Single point of failure", "Layered — ML can fail without affecting scoring"],
        ["SOC Workflow", "Static CSV/SQL", "Basic charts", "Live WebSocket Feed, Kanban & Graph Explorer"],
        ["Demonstrated Precision", "Varies (~60-70%)", "Varies (~80-85%)", "100.0% Precision (0 False Positives)"],
        ["Attack Simulator", "None", "None", "1-Click Scenario Injector + Continuous Auto-Flow"]
    ]
    for i, row in enumerate(matrix_rows):
        for j, val in enumerate(row):
            cell = table.cell(i+1, j)
            cell.fill.solid()
            if j == 3:
                cell.fill.fore_color.rgb = RGBColor(20, 35, 60)
            else:
                cell.fill.fore_color.rgb = RGBColor(10, 15, 26) if i % 2 == 0 else RGBColor(15, 23, 42)

            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.size = Pt(10)
            if j == 3:
                p.font.bold = True
                p.font.color.rgb = NEON_BLUE
            else:
                p.font.color.rgb = SLATE_LIGHT

    # ----------------------------------------------------
    # SLIDE 15: Future Vision & Product Roadmap
    # ----------------------------------------------------
    s15 = prs.slides.add_slide(slide_layout)
    apply_dark_bg(s15)
    add_header(s15, "The Horizon: Next-Gen Autonomous Fraud Defense", "Scaling PayTelemetry Into a Cross-Institutional Intelligence Network", "15")

    phases = [
        {"phase": "PHASE 1 (Near-Term)", "title": "On-Device Pre-Check Scoring", "col": NEON_BLUE, "desc": "Running lightweight exported models directly on mobile devices for zero-network pre-checks."},
        {"phase": "PHASE 2 (Mid-Term)", "title": "Privacy-Preserving Cross-Bank Graph", "col": EMERALD, "desc": "Cross-institution graph collaboration for mule ring detection without sharing raw customer PII using Homomorphic Encryption."},
        {"phase": "PHASE 3 (Long-Term)", "title": "LLM-Assisted SOC Copilots", "col": AMBER, "desc": "Copilots that draft case summaries and recommend next actions — always with a human investigator in the loop, never autonomous action."}
    ]
    for i, ph in enumerate(phases):
        bx = s15.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + i*4.0), Inches(1.6), Inches(3.733), Inches(5.1))
        bx.fill.solid()
        bx.fill.fore_color.rgb = CARD_DARK
        bx.line.color.rgb = CARD_BORDER
        tf = bx.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = ph["phase"]
        p.font.size = Pt(10.5)
        p.font.bold = True
        p.font.color.rgb = ph["col"]

        p_t = tf.add_paragraph()
        p_t.text = ph["title"]
        p_t.font.size = Pt(14)
        p_t.font.bold = True
        p_t.font.color.rgb = WHITE

        tf.add_paragraph()
        p_d = tf.add_paragraph()
        p_d.text = ph["desc"]
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = SLATE_LIGHT

    # ----------------------------------------------------
    # SLIDE 16: Thank You & Open Q&A (Final Slide)
    # ----------------------------------------------------
    s16 = prs.slides.add_slide(slide_layout)
    apply_dark_bg(s16)

    b_s16 = s16.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(1.0), Inches(10.933), Inches(5.5))
    b_s16.fill.solid()
    b_s16.fill.fore_color.rgb = CARD_DARK
    b_s16.line.color.rgb = ROYAL_BLUE
    b_s16.line.width = Pt(2)
    tf16 = b_s16.text_frame
    tf16.word_wrap = True

    p = tf16.paragraphs[0]
    p.text = "HEXAWARE MAVERICKS HACKATHON 2026"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = NEON_BLUE
    p.alignment = PP_ALIGN.CENTER

    p_t = tf16.add_paragraph()
    p_t.text = "Thank You!"
    p_t.font.size = Pt(36)
    p_t.font.bold = True
    p_t.font.color.rgb = WHITE
    p_t.alignment = PP_ALIGN.CENTER

    p_sub = tf16.add_paragraph()
    p_sub.text = "Team Chronoforge is Ready for Your Questions & Live System Demonstration."
    p_sub.font.size = Pt(15)
    p_sub.font.color.rgb = SLATE_MUTED
    p_sub.alignment = PP_ALIGN.CENTER

    tf16.add_paragraph()
    p_q = tf16.add_paragraph()
    p_q.text = "\"Security should protect trust, not destroy user experience. PayTelemetry makes real-time payments safer, smarter, and frictionless.\""
    p_q.font.size = Pt(13)
    p_q.font.bold = True
    p_q.font.color.rgb = EMERALD
    p_q.alignment = PP_ALIGN.CENTER

    tf16.add_paragraph()
    p_links = tf16.add_paragraph()
    p_links.text = "📱 Consumer UPI App: http://localhost:5173  |  🛡️ Investigator Console: http://localhost:5173"
    p_links.font.size = Pt(11)
    p_links.font.color.rgb = SLATE_LIGHT
    p_links.alignment = PP_ALIGN.CENTER

    p_repo = tf16.add_paragraph()
    p_repo.text = "📦 GitHub Repository: https://github.com/Balaji-9941/HACKATHON-PROJECT"
    p_repo.font.size = Pt(11)
    p_repo.font.color.rgb = NEON_BLUE
    p_repo.alignment = PP_ALIGN.CENTER

    # Save presentation
    output_path = os.path.abspath(os.path.join(os.path.dirname(__file__), 'PayTelemetry_Technical_Judging_Deck_Chronoforge.pptx'))
    prs.save(output_path)
    print(f"[PPT] 16-Slide Deck successfully created at: {output_path}")

if __name__ == '__main__':
    build_deck()
